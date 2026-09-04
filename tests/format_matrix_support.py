"""Fixture builders for the format-compatibility matrix acceptance test.

Every builder produces the smallest lawful fixture that exercises the real
ingestion path for one matrix family.  Patterns mirror the dedicated
per-format test suites so both agree on what a valid input is.
"""

from __future__ import annotations

import math
import struct
import wave
from io import BytesIO
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZIP_STORED, ZipFile

from docx import Document
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Emu
from reportlab.pdfgen.canvas import Canvas
from xlwt import Workbook as LegacyWorkbook

MARKER = "ARCHIV-MATRIX-MARKER"

OFFICE = "urn:oasis:names:tc:opendocument:xmlns:office:1.0"
TEXT = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
TABLE = "urn:oasis:names:tc:opendocument:xmlns:table:1.0"
DRAW = "urn:oasis:names:tc:opendocument:xmlns:drawing:1.0"
PRESENTATION = "urn:oasis:names:tc:opendocument:xmlns:presentation:1.0"
DATABASE = "urn:oasis:names:tc:opendocument:xmlns:database:1.0"
MANIFEST = "urn:oasis:names:tc:opendocument:xmlns:manifest:1.0"
MATHML = "http://www.w3.org/1998/Math/MathML"
XLINK = "http://www.w3.org/1999/xlink"

FREESECT = 0xFFFFFFFF
ENDOFCHAIN = 0xFFFFFFFE
FATSECT = 0xFFFFFFFD


def build_text() -> bytes:
    return f"Archiv format matrix plain-text fixture.\n{MARKER} line two.\n".encode()


def build_pdf() -> bytes:
    raw = BytesIO()
    canvas = Canvas(raw, pagesize=(612, 792), invariant=1, pageCompression=0)
    canvas.setTitle("Archiv matrix PDF fixture")
    canvas.setFont("Helvetica", 14)
    canvas.drawString(72, 720, "Archiv matrix PDF fixture")
    canvas.drawString(72, 692, MARKER)
    canvas.showPage()
    canvas.save()
    return raw.getvalue()


def build_blank_pdf() -> bytes:
    raw = BytesIO()
    canvas = Canvas(raw, pagesize=(612, 792), invariant=1, pageCompression=0)
    canvas.setTitle("Archiv matrix blank PDF fixture")
    canvas.showPage()
    canvas.save()
    return raw.getvalue()


def build_docx() -> bytes:
    document = Document()
    document.add_paragraph("Archiv matrix DOCX fixture")
    document.add_paragraph(MARKER)
    raw = BytesIO()
    document.save(raw)
    return raw.getvalue()


def build_rtf() -> bytes:
    return (
        r"{\rtf1\ansi\ansicpg1252\deff0"
        r"{\fonttbl{\f0 Calibri;}}"
        r"\f0 Archiv matrix RTF fixture\par " + MARKER + r"\par}"
    ).encode("ascii")


def build_xlsx() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Evidence"
    sheet["A1"] = "Archiv matrix XLSX fixture"
    sheet["B2"] = MARKER
    raw = BytesIO()
    workbook.save(raw)
    return raw.getvalue()


def build_xls() -> bytes:
    workbook = LegacyWorkbook(encoding="utf-8")
    sheet = workbook.add_sheet("Evidence")
    sheet.write(0, 0, "Archiv matrix XLS fixture")
    sheet.write(1, 1, MARKER)
    raw = BytesIO()
    workbook.save(raw)
    return raw.getvalue()


def build_pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Emu(457200), Emu(457200), Emu(8229600), Emu(914400))
    title.text_frame.text = "Archiv matrix PPTX fixture"
    marker = slide.shapes.add_textbox(Emu(457200), Emu(1828800), Emu(8229600), Emu(914400))
    marker.text_frame.text = MARKER
    raw = BytesIO()
    presentation.save(raw)
    return raw.getvalue()


def _odf_manifest(mimetype: str) -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8"?>'
        f'<manifest:manifest xmlns:manifest="{MANIFEST}">'
        f'<manifest:file-entry manifest:full-path="/" manifest:media-type="{mimetype}"/>'
        '<manifest:file-entry manifest:full-path="content.xml" manifest:media-type="text/xml"/>'
        "</manifest:manifest>"
    )


def _odf_package(path: Path, mimetype: str, content: str) -> None:
    with ZipFile(path, "w") as archive:
        archive.writestr("mimetype", mimetype, compress_type=ZIP_STORED)
        archive.writestr("content.xml", content, compress_type=ZIP_DEFLATED)
        archive.writestr(
            "META-INF/manifest.xml", _odf_manifest(mimetype), compress_type=ZIP_DEFLATED
        )


def _odf_package_document(body: str) -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8"?>'
        f'<office:document-content xmlns:office="{OFFICE}" xmlns:text="{TEXT}" '
        f'xmlns:table="{TABLE}" xmlns:draw="{DRAW}" xmlns:presentation="{PRESENTATION}">'
        f"<office:body>{body}</office:body></office:document-content>"
    )


def _odf_flat_document(mimetype: str, body: str) -> str:
    return (
        '<?xml version="1.0" encoding="UTF-8"?>'
        f'<office:document xmlns:office="{OFFICE}" xmlns:text="{TEXT}" '
        f'xmlns:table="{TABLE}" xmlns:draw="{DRAW}" xmlns:presentation="{PRESENTATION}" '
        f'office:mimetype="{mimetype}">'
        f"<office:body>{body}</office:body></office:document>"
    )


def _odf_body(family: str) -> str:
    if family == "text":
        return f"<office:text><text:p>{MARKER}</text:p></office:text>"
    if family == "spreadsheet":
        return (
            '<office:spreadsheet><table:table table:name="Sheet">'
            "<table:table-row><table:table-cell>"
            f"<text:p>{MARKER}</text:p>"
            "</table:table-cell></table:table-row></table:table></office:spreadsheet>"
        )
    container = "presentation" if family == "presentation" else "drawing"
    return (
        f"<office:{container}><draw:page><draw:frame><draw:text-box>"
        f"<text:p>{MARKER}</text:p>"
        f"</draw:text-box></draw:frame></draw:page></office:{container}>"
    )


_ODF_PACKAGES: dict[str, tuple[str, str]] = {
    "odt": ("application/vnd.oasis.opendocument.text", "text"),
    "ott": ("application/vnd.oasis.opendocument.text-template", "text"),
    "odm": ("application/vnd.oasis.opendocument.text-master", "text"),
    "otm": ("application/vnd.oasis.opendocument.text-master-template", "text"),
    "ods": ("application/vnd.oasis.opendocument.spreadsheet", "spreadsheet"),
    "ots": ("application/vnd.oasis.opendocument.spreadsheet-template", "spreadsheet"),
    "odp": ("application/vnd.oasis.opendocument.presentation", "presentation"),
    "otp": ("application/vnd.oasis.opendocument.presentation-template", "presentation"),
    "odg": ("application/vnd.oasis.opendocument.graphics", "drawing"),
    "otg": ("application/vnd.oasis.opendocument.graphics-template", "drawing"),
}

_ODF_FLAT: dict[str, tuple[str, str]] = {
    "fodt": ("application/vnd.oasis.opendocument.text", "text"),
    "fods": ("application/vnd.oasis.opendocument.spreadsheet", "spreadsheet"),
    "fodp": ("application/vnd.oasis.opendocument.presentation", "presentation"),
    "fodg": ("application/vnd.oasis.opendocument.graphics", "drawing"),
}


def build_odf_package(path: Path, suffix: str) -> None:
    mimetype, family = _ODF_PACKAGES[suffix]
    _odf_package(path, mimetype, _odf_package_document(_odf_body(family)))


def build_odf_flat(path: Path, suffix: str) -> None:
    mimetype, family = _ODF_FLAT[suffix]
    path.write_text(_odf_flat_document(mimetype, _odf_body(family)), encoding="utf-8")


def build_odf_formula(path: Path) -> None:
    content = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        f'<math:math xmlns:math="{MATHML}"><math:mrow>'
        "<math:mi>x</math:mi><math:mo>+</math:mo><math:mn>1</math:mn>"
        "</math:mrow></math:math>"
    )
    _odf_package(path, "application/vnd.oasis.opendocument.formula", content)


def build_odb(path: Path) -> None:
    names = ["ArchivMatrixEvidence"]
    representations = "".join(f'<db:table-representation db:name="{name}"/>' for name in names)
    body = (
        "<office:database>"
        "<db:data-source><db:connection-data>"
        '<db:connection-resource xlink:href="sdbc:embedded:firebird"/>'
        "</db:connection-data></db:data-source>"
        f"<db:table-representations>{representations}</db:table-representations>"
        "</office:database>"
    )
    content = (
        '<?xml version="1.0" encoding="UTF-8"?>'
        f'<office:document-content xmlns:office="{OFFICE}" xmlns:db="{DATABASE}" '
        f'xmlns:xlink="{XLINK}" office:version="1.2">'
        f"<office:body>{body}</office:body></office:document-content>"
    )
    with ZipFile(path, "w") as archive:
        archive.writestr(
            "mimetype", "application/vnd.oasis.opendocument.base", compress_type=ZIP_STORED
        )
        archive.writestr("content.xml", content, compress_type=ZIP_DEFLATED)
        archive.writestr(
            "META-INF/manifest.xml",
            _odf_manifest("application/vnd.oasis.opendocument.base"),
            compress_type=ZIP_DEFLATED,
        )


# The legacy InPage container builder mirrors the InPage ingestion test so both
# suites construct identical synthetic fixtures.
def _cfb_entry(
    name: str,
    *,
    object_type: int,
    right: int = FREESECT,
    child: int = FREESECT,
    start: int = ENDOFCHAIN,
    size: int = 0,
) -> bytes:
    encoded = (name + "\x00").encode("utf-16le")
    raw = bytearray(128)
    raw[: len(encoded)] = encoded
    struct.pack_into("<H", raw, 64, len(encoded))
    raw[66] = object_type
    raw[67] = 1
    struct.pack_into("<III", raw, 68, FREESECT, right, child)
    struct.pack_into("<I", raw, 116, start)
    struct.pack_into("<Q", raw, 120, size)
    return bytes(raw)


def build_inpage300(path: Path) -> None:
    payload = f"{MARKER} اردو\r\nدوسری سطر".encode("utf-16le")
    payload = payload + b"\x00" * max(0, 4096 - len(payload))
    document = b"D" * 4096
    doc_sectors = 8
    native_sectors = (len(payload) + 511) // 512
    total_sectors = 2 + doc_sectors + native_sectors
    header = bytearray(512)
    header[:8] = bytes.fromhex("D0CF11E0A1B11AE1")
    struct.pack_into("<HHHH", header, 24, 0x003E, 3, 0xFFFE, 9)
    struct.pack_into("<H", header, 32, 6)
    struct.pack_into("<I", header, 44, 1)
    struct.pack_into("<I", header, 48, 0)
    struct.pack_into("<I", header, 56, 4096)
    struct.pack_into("<I", header, 60, ENDOFCHAIN)
    struct.pack_into("<I", header, 68, ENDOFCHAIN)
    struct.pack_into("<109I", header, 76, 1, *([FREESECT] * 108))

    name = "InPage300"
    directory = bytearray(512)
    directory[:128] = _cfb_entry("Root Entry", object_type=5, child=1)
    directory[128:256] = _cfb_entry(
        "DocumentInfo", object_type=2, right=2, start=2, size=len(document)
    )
    directory[256:384] = _cfb_entry(
        name,
        object_type=2,
        start=2 + doc_sectors,
        size=len(payload),
    )

    fat = [FREESECT] * 128
    fat[0] = ENDOFCHAIN
    fat[1] = FATSECT
    for sector in range(2, 2 + doc_sectors):
        fat[sector] = sector + 1 if sector < 1 + doc_sectors else ENDOFCHAIN
    start = 2 + doc_sectors
    for sector in range(start, start + native_sectors):
        fat[sector] = sector + 1 if sector < start + native_sectors - 1 else ENDOFCHAIN
    body = bytearray(directory + struct.pack("<128I", *fat) + document + payload)
    assert len(body) == total_sectors * 512
    path.write_bytes(bytes(header + body))


def _cfb_container(streams: dict[str, bytes]) -> bytes:
    """Build a minimal multi-stream CFB container (root + up to 3 sibling streams)."""

    assert 1 <= len(streams) <= 3
    padded = {
        name: (data if len(data) >= 4096 else data + b"\x00" * (4096 - len(data)))
        for name, data in streams.items()
    }
    header = bytearray(512)
    header[:8] = bytes.fromhex("D0CF11E0A1B11AE1")
    struct.pack_into("<HHHH", header, 24, 0x003E, 3, 0xFFFE, 9)
    struct.pack_into("<H", header, 32, 6)
    struct.pack_into("<I", header, 44, 1)
    struct.pack_into("<I", header, 48, 0)
    struct.pack_into("<I", header, 56, 4096)
    struct.pack_into("<I", header, 60, ENDOFCHAIN)
    struct.pack_into("<I", header, 68, ENDOFCHAIN)
    struct.pack_into("<109I", header, 76, 1, *([FREESECT] * 108))

    directory = bytearray(512)
    directory[:128] = _cfb_entry("Root Entry", object_type=5, child=1)
    sector = 2
    body = bytearray()
    names = list(padded)
    for index, name in enumerate(names, start=1):
        data = padded[name]
        sectors = len(data) // 512
        right = index + 1 if index < len(names) else FREESECT
        # Declared size must be >= the mini-stream cutoff (4096) or the CFB reader
        # expects mini-FAT allocation instead of the regular FAT chain below.
        directory[index * 128 : (index + 1) * 128] = _cfb_entry(
            name, object_type=2, right=right, start=sector, size=len(data)
        )
        body.extend(data)
        sector += sectors

    fat_len = ((sector + 127) // 128) * 128
    fat = [FREESECT] * fat_len
    fat[0] = ENDOFCHAIN
    fat[1] = FATSECT
    cursor = 2
    for name in names:
        sectors = len(padded[name]) // 512
        for offset in range(sectors):
            fat[cursor + offset] = cursor + offset + 1 if offset < sectors - 1 else ENDOFCHAIN
        cursor += sectors

    return bytes(header + directory + struct.pack(f"<{fat_len}I", *fat) + body)


def build_doc() -> bytes:
    """A minimal but spec-accurate Word 97 (.doc) fixture: FIB + Clx + one text piece."""

    text = f"Archiv matrix DOC fixture {MARKER}\r"
    text_bytes = text.encode("utf-16-le")

    word_document = bytearray(1024)
    struct.pack_into("<H", word_document, 0, 0xA5EC)  # wIdent
    struct.pack_into("<H", word_document, 2, 0x00C1)  # nFib = Word 97
    struct.pack_into("<H", word_document, 32, 14)  # csw
    struct.pack_into("<H", word_document, 62, 22)  # cslw
    struct.pack_into("<H", word_document, 152, 34)  # cbRgFcLcb
    text_offset = 512
    word_document[text_offset : text_offset + len(text_bytes)] = text_bytes

    cp_values = struct.pack("<II", 0, len(text))  # one piece: CP[0]=0, CP[1]=char count
    pcd = struct.pack("<HIH", 0, text_offset, 0)  # flags, FcCompressed(fCompressed=0), prm
    plc_pcd = cp_values + pcd
    clx = bytes([0x02]) + struct.pack("<I", len(plc_pcd)) + plc_pcd

    fc_clx_offset = 154 + 33 * 8
    struct.pack_into("<II", word_document, fc_clx_offset, 0, len(clx))  # fcClx=0 in "0Table"

    return _cfb_container({"WordDocument": bytes(word_document), "0Table": clx})


def build_ppt() -> bytes:
    """A minimal PowerPoint 97 (.ppt) fixture: one slide with one text run."""

    def record(rec_type: int, payload: bytes, *, container: bool = False) -> bytes:
        ver_instance = 0x000F if container else 0x0000
        return struct.pack("<HHI", ver_instance, rec_type, len(payload)) + payload

    text = f"Archiv matrix PPT fixture {MARKER}"
    text_atom = record(0x0FA0, text.encode("utf-16-le"))
    slide = record(0x03EE, text_atom, container=True)
    document = record(0x03E8, slide, container=True)

    return _cfb_container({"PowerPoint Document": document})


def build_png() -> bytes:
    image = Image.new("L", (320, 40), 255)
    raw = BytesIO()
    image.save(raw, format="PNG")
    return raw.getvalue()


def build_jpeg() -> bytes:
    image = Image.new("RGB", (320, 40), (255, 255, 255))
    raw = BytesIO()
    image.save(raw, format="JPEG")
    return raw.getvalue()


def build_gif() -> bytes:
    image = Image.new("P", (320, 40), 0)
    raw = BytesIO()
    image.save(raw, format="GIF")
    return raw.getvalue()


def build_bmp() -> bytes:
    image = Image.new("RGB", (320, 40), (255, 255, 255))
    raw = BytesIO()
    image.save(raw, format="BMP")
    return raw.getvalue()


def build_tiff() -> bytes:
    image = Image.new("RGB", (320, 40), (255, 255, 255))
    raw = BytesIO()
    image.save(raw, format="TIFF")
    return raw.getvalue()


def build_webp() -> bytes:
    image = Image.new("RGB", (320, 40), (255, 255, 255))
    raw = BytesIO()
    image.save(raw, format="WEBP")
    return raw.getvalue()


def build_wav() -> bytes:
    raw = BytesIO()
    with wave.open(raw, "wb") as audio:
        audio.setnchannels(1)
        audio.setsampwidth(2)
        audio.setframerate(8000)
        frames = bytearray()
        for index in range(800):
            value = int(1000 * math.sin(2 * math.pi * 440 * index / 8000))
            frames.extend(struct.pack("<h", value))
        audio.writeframes(bytes(frames))
    return raw.getvalue()


def build_svg() -> bytes:
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">\n'
        f"  <title>Archiv matrix SVG fixture</title>\n"
        f"  <desc>SVG format compatibility test</desc>\n"
        f'  <text x="10" y="20">{MARKER}</text>\n'
        f"</svg>\n"
    ).encode()
