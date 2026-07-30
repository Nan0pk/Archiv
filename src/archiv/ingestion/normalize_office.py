"""Spreadsheet and presentation normalization."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation

from archiv.contracts import NormalizedDocument, NormalizedSegment, NormalizedTable


def normalize_xlsx(
    path: Path,
    digest: str,
    *,
    source_name: str,
    media_type: str,
) -> NormalizedDocument:
    workbook = load_workbook(BytesIO(path.read_bytes()), read_only=True, data_only=False)
    segments: list[NormalizedSegment] = []
    tables: list[NormalizedTable] = []
    for worksheet in workbook.worksheets:
        rows: list[list[object | None]] = []
        for row in worksheet.iter_rows():
            values: list[object | None] = []
            for cell in row:
                value = cell.value
                values.append(value)
                if value is not None:
                    segments.append(
                        NormalizedSegment(
                            locator={"sheet": worksheet.title, "cell": cell.coordinate},
                            text=str(value),
                        )
                    )
            if any(value is not None for value in values):
                rows.append(values)
        if rows:
            tables.append(NormalizedTable(locator={"sheet": worksheet.title}, rows=rows))
    return NormalizedDocument(
        object_sha256=digest,
        media_type=media_type,
        kind="xlsx",
        source_name=source_name,
        segments=segments,
        tables=tables,
        metadata={"sheets": workbook.sheetnames},
    )


def normalize_pptx(
    path: Path,
    digest: str,
    *,
    source_name: str,
    media_type: str,
) -> NormalizedDocument:
    presentation = Presentation(str(path))
    segments: list[NormalizedSegment] = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        for shape_number, shape in enumerate(slide.shapes, 1):
            text = getattr(shape, "text", "")
            if text:
                segments.append(
                    NormalizedSegment(
                        locator={"slide": slide_number, "shape": shape_number},
                        text=text,
                    )
                )
    return NormalizedDocument(
        object_sha256=digest,
        media_type=media_type,
        kind="pptx",
        source_name=source_name,
        segments=segments,
        metadata={"slides": len(presentation.slides)},
    )
