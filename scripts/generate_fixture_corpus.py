#!/usr/bin/env python3
"""Generate Archiv's deterministic, public-safe representative fixture corpus."""

from __future__ import annotations

import argparse
import hashlib
import json
import math
import re
import struct
from datetime import UTC, datetime
from io import BytesIO
from pathlib import Path
from typing import Final
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from reportlab.pdfgen.canvas import Canvas

FIXED_TIME: Final = (1980, 1, 1, 0, 0, 0)
FIXED_DATETIME: Final = datetime(2026, 1, 1, tzinfo=UTC)
DEFAULT_OUTPUT: Final = Path("build/fixtures/representative-corpus")

FIXTURES: Final[dict[str, dict[str, object]]] = {
    "plain-text.txt": {
        "media_type": "text/plain",
        "marker": "ARCHIV-TEXT-MARKER-2026",
        "location": {"line": 3},
        "expected_valid": True,
    },
    "report.pdf": {
        "media_type": "application/pdf",
        "marker": "ARCHIV-PDF-MARKER-2026",
        "location": {"page": 1},
        "expected_valid": True,
    },
    "document.docx": {
        "media_type": (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ),
        "marker": "ARCHIV-DOCX-MARKER-2026",
        "location": {"paragraph": 2},
        "expected_valid": True,
    },
    "workbook.xlsx": {
        "media_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "marker": "ARCHIV-XLSX-MARKER-2026",
        "location": {"sheet": "Evidence", "cell": "B2"},
        "expected_valid": True,
    },
    "presentation.pptx": {
        "media_type": (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ),
        "marker": "ARCHIV-PPTX-MARKER-2026",
        "location": {"slide": 1, "shape": 2},
        "expected_valid": True,
    },
    "scanned-page.png": {
        "media_type": "image/png",
        "marker": "ARCHIV-IMAGE-MARKER-2026",
        "location": {"page": 1, "bounding_box": [20, 20, 246, 31]},
        "expected_valid": True,
    },
    "sample.wav": {
        "media_type": "audio/wav",
        "marker": "ARCHIV-AUDIO-MARKER-2026",
        "location": {
            "metadata_chunk": "LIST/INFO/ICMT",
            "tone_segments_ms": [[0, 250], [300, 550], [600, 850]],
        },
        "expected_valid": True,
    },
    "malformed/truncated.pdf": {
        "media_type": "application/pdf",
        "marker": None,
        "location": {"expected": "reject"},
        "expected_valid": False,
    },
    "malformed/not-a-docx.docx": {
        "media_type": (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ),
        "marker": None,
        "location": {"expected": "reject"},
        "expected_valid": False,
    },
    "malformed/corrupt.wav": {
        "media_type": "audio/wav",
        "marker": None,
        "location": {"expected": "reject"},
        "expected_valid": False,
    },
}


def _normalize_zip(raw: bytes) -> bytes:
    source = BytesIO(raw)
    target = BytesIO()
    with ZipFile(source) as input_zip, ZipFile(target, "w") as output_zip:
        for name in sorted(input_zip.namelist()):
            info = ZipInfo(name, FIXED_TIME)
            info.compress_type = ZIP_DEFLATED
            info.create_system = 3
            info.external_attr = 0o100644 << 16
            content = input_zip.read(name)
            if name == "docProps/core.xml":
                content = re.sub(
                    rb"<dcterms:(created|modified)([^>]*)>[^<]+</dcterms:\1>",
                    lambda match: (
                        b"<dcterms:"
                        + match.group(1)
                        + match.group(2)
                        + b">2026-01-01T00:00:00Z</dcterms:"
                        + match.group(1)
                        + b">"
                    ),
                    content,
                )
            output_zip.writestr(info, content)
    return target.getvalue()


def _docx() -> bytes:
    document = Document()
    document.core_properties.title = "Archiv DOCX Fixture"
    document.core_properties.author = "Archiv fixture generator"
    document.core_properties.created = FIXED_DATETIME
    document.core_properties.modified = FIXED_DATETIME
    document.add_paragraph("Archiv DOCX Fixture")
    document.add_paragraph("ARCHIV-DOCX-MARKER-2026")
    document.add_paragraph("Section: Findings; paragraph location: 2")
    raw = BytesIO()
    document.save(raw)
    return _normalize_zip(raw.getvalue())


def _xlsx() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Evidence"
    sheet["A1"] = "Archiv XLSX Fixture"
    sheet["B2"] = "ARCHIV-XLSX-MARKER-2026"
    sheet["C3"] = 42.25
    workbook.properties.title = "Archiv XLSX Fixture"
    workbook.properties.creator = "Archiv fixture generator"
    workbook.properties.created = FIXED_DATETIME
    workbook.properties.modified = FIXED_DATETIME
    raw = BytesIO()
    workbook.save(raw)
    return _normalize_zip(raw.getvalue())


def _pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(457200, 457200, 8229600, 914400)
    title.text_frame.text = "Archiv PPTX Fixture"
    marker = slide.shapes.add_textbox(457200, 1828800, 8229600, 914400)
    marker.text_frame.text = "ARCHIV-PPTX-MARKER-2026"
    presentation.core_properties.title = "Archiv PPTX Fixture"
    presentation.core_properties.author = "Archiv fixture generator"
    presentation.core_properties.created = FIXED_DATETIME
    presentation.core_properties.modified = FIXED_DATETIME
    raw = BytesIO()
    presentation.save(raw)
    return _normalize_zip(raw.getvalue())


def _pdf() -> bytes:
    raw = BytesIO()
    canvas = Canvas(
        raw,
        pagesize=(612, 792),
        invariant=1,
        pageCompression=0,
        pdfVersion=(1, 4),
    )
    canvas.setTitle("Archiv PDF Fixture")
    canvas.setAuthor("Archiv fixture generator")
    canvas.setFont("Helvetica", 18)
    canvas.drawString(72, 720, "Archiv PDF Fixture")
    canvas.drawString(72, 692, "ARCHIV-PDF-MARKER-2026")
    canvas.drawString(72, 664, "Page location: 1")
    canvas.showPage()
    canvas.save()
    return raw.getvalue()


def _png() -> bytes:
    image = Image.new("L", (640, 80), 255)
    draw = ImageDraw.Draw(image)
    draw.text(
        (20, 20),
        "ARCHIV-IMAGE-MARKER-2026",
        fill=0,
        font=ImageFont.load_default(),
    )
    raw = BytesIO()
    image.save(raw, format="PNG", optimize=False, compress_level=9)
    return raw.getvalue()


def _chunk(kind: bytes, payload: bytes) -> bytes:
    return (
        kind
        + struct.pack("<I", len(payload))
        + payload
        + (b"\0" if len(payload) % 2 else b"")
    )


def _wav() -> bytes:
    sample_rate = 8000
    tones = (
        (440.0, 0.25),
        (0.0, 0.05),
        (554.37, 0.25),
        (0.0, 0.05),
        (659.25, 0.25),
    )
    samples: list[int] = []
    for frequency, duration in tones:
        for index in range(round(sample_rate * duration)):
            value = (
                0
                if frequency == 0
                else int(
                    9000
                    * math.sin(2 * math.pi * frequency * index / sample_rate)
                )
            )
            samples.append(value)
    pcm = struct.pack("<" + "h" * len(samples), *samples)
    fmt = struct.pack("<HHIIHH", 1, 1, sample_rate, sample_rate * 2, 2, 16)
    info = b"INFO" + _chunk(b"ICMT", b"ARCHIV-AUDIO-MARKER-2026\0")
    chunks = _chunk(b"fmt ", fmt) + _chunk(b"data", pcm) + _chunk(b"LIST", info)
    return b"RIFF" + struct.pack("<I", len(chunks) + 4) + b"WAVE" + chunks


def _provenance() -> bytes:
    return b"""# Fixture provenance and licensing

All files are generated from synthetic constants by `scripts/generate_fixture_corpus.py`.
No personal, organisational, customer, classified, scraped, or proprietary material is used.

The fixtures have no separate licence and follow the repository's current default terms.
At introduction time Archiv was public but had not selected an open-source licence, so no
reuse rights should be inferred solely from public visibility.

The generator pins its document/image libraries in the development dependency set, fixes
Office metadata, normalizes ZIP ordering/timestamps/permissions, uses invariant PDF output,
and fixes image pixels and WAV samples. LibreOffice headless conversion was used as an
additional compatibility check for the generated DOCX, XLSX, PPTX, and PDF.
"""


def build_files() -> dict[str, bytes]:
    expected = {
        "schema_version": 1,
        "fixtures": FIXTURES,
        "spreadsheet_values": {"workbook.xlsx": {"Evidence!C3": 42.25}},
        "audio_note": (
            "WAV tones test deterministic timing; the marker is stored in "
            "LIST/INFO/ICMT, not spoken."
        ),
    }
    return {
        "plain-text.txt": (
            b"Archiv plain-text fixture\n"
            b"Location test follows\n"
            b"ARCHIV-TEXT-MARKER-2026\n"
            b"End of fixture\n"
        ),
        "report.pdf": _pdf(),
        "document.docx": _docx(),
        "workbook.xlsx": _xlsx(),
        "presentation.pptx": _pptx(),
        "scanned-page.png": _png(),
        "sample.wav": _wav(),
        "malformed/truncated.pdf": (
            b"%PDF-1.7\n1 0 obj\n<< /Type /Catalog >>\nendobj\n"
        ),
        "malformed/not-a-docx.docx": (
            b"ARCHIV malformed DOCX fixture: intentionally not a ZIP package.\n"
        ),
        "malformed/corrupt.wav": b"RIFF\x08\0\0\0NOTWAVE!",
        "expected.json": (
            json.dumps(expected, indent=2, sort_keys=True) + "\n"
        ).encode(),
        "PROVENANCE.md": _provenance(),
    }


def generate(output: Path) -> None:
    files = build_files()
    output.mkdir(parents=True, exist_ok=True)
    for relative, content in files.items():
        path = output / relative
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_bytes(content)

    entries = [
        {
            "path": path,
            "bytes": len(content),
            "sha256": hashlib.sha256(content).hexdigest(),
            "media_type": FIXTURES.get(path, {}).get(
                "media_type",
                "application/json" if path.endswith(".json") else "text/markdown",
            ),
            "marker": FIXTURES.get(path, {}).get("marker"),
            "license": "repository-default; no separate fixture licence",
        }
        for path, content in sorted(files.items())
    ]
    manifest = {
        "schema_version": 1,
        "generator_version": 1,
        "generator": "scripts/generate_fixture_corpus.py",
        "entries": entries,
    }
    (output / "manifest.json").write_text(
        json.dumps(manifest, indent=2, sort_keys=True) + "\n",
        encoding="utf-8",
    )


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    generate(parser.parse_args().output)


if __name__ == "__main__":
    main()
