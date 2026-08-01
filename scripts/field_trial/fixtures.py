"""Deterministic public fixtures and fake loopback model."""

from __future__ import annotations

import json
import re
import threading
from collections.abc import Mapping, Sequence
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from io import BytesIO
from pathlib import Path
from typing import cast

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen.canvas import Canvas

from field_trial.common import (
    BenchmarkError,
    FIXED_DATETIME,
    SCHEMA_VERSION,
    _normalize_zip,
    sha256_file,
)


def _write_fixture(path: Path, source: Mapping[str, object]) -> None:
    file_format = str(source["format"])
    if file_format in {"text", "markdown"}:
        path.write_text(str(source["content"]), encoding="utf-8", newline="\n")
        return
    if file_format == "docx":
        document = Document()
        document.core_properties.created = FIXED_DATETIME
        document.core_properties.modified = FIXED_DATETIME
        for paragraph in cast(Sequence[object], source["paragraphs"]):
            document.add_paragraph(str(paragraph))
        raw = BytesIO()
        document.save(raw)
        path.write_bytes(_normalize_zip(raw.getvalue()))
        return
    if file_format == "xlsx":
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Evidence"
        for row in cast(Sequence[Sequence[object]], source["rows"]):
            sheet.append(list(row))
        workbook.properties.created = FIXED_DATETIME
        workbook.properties.modified = FIXED_DATETIME
        raw = BytesIO()
        workbook.save(raw)
        path.write_bytes(_normalize_zip(raw.getvalue()))
        return
    if file_format == "pptx":
        presentation = Presentation()
        for slide_spec in cast(Sequence[Mapping[str, object]], source["slides"]):
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            if slide.shapes.title is not None:
                slide.shapes.title.text = str(slide_spec["title"])
            frame = slide.placeholders[1].text_frame
            frame.clear()
            for index, bullet in enumerate(cast(Sequence[object], slide_spec["bullets"])):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = str(bullet)
        presentation.core_properties.created = FIXED_DATETIME
        presentation.core_properties.modified = FIXED_DATETIME
        raw = BytesIO()
        presentation.save(raw)
        path.write_bytes(_normalize_zip(raw.getvalue()))
        return
    if file_format == "pdf":
        raw = BytesIO()
        canvas = Canvas(raw, pagesize=(612, 792), invariant=1, pageCompression=0)
        y = 740
        for paragraph in cast(Sequence[object], source["paragraphs"]):
            canvas.drawString(72, y, str(paragraph))
            y -= 28
        canvas.showPage()
        canvas.save()
        path.write_bytes(raw.getvalue())
        return
    raise BenchmarkError(f"unsupported fixture format: {file_format}")


def generate_public_corpus(
    benchmark: Mapping[str, object], output: Path
) -> list[dict[str, object]]:
    output.mkdir(parents=True, exist_ok=True)
    manifest: list[dict[str, object]] = []
    for source in cast(Sequence[Mapping[str, object]], benchmark["corpus"]):
        path = output / str(source["filename"])
        _write_fixture(path, source)
        manifest.append(
            {
                "source_id": source["id"],
                "filename": path.name,
                "format": source["format"],
                "sha256": sha256_file(path),
                "size_bytes": path.stat().st_size,
            }
        )
    return sorted(manifest, key=lambda item: str(item["source_id"]))


def _source_maps(
    benchmark: Mapping[str, object],
) -> tuple[dict[str, str], dict[str, str]]:
    by_id: dict[str, str] = {}
    by_filename: dict[str, str] = {}
    for source in cast(Sequence[Mapping[str, object]], benchmark["corpus"]):
        source_id = str(source["id"])
        filename = str(source["filename"])
        by_id[source_id] = filename
        by_filename[filename] = source_id
    return by_id, by_filename


def _fake_response(benchmark: Mapping[str, object], prompt: str) -> str:
    if prompt.strip() == "Respond with the single word PONG.":
        return "PONG"
    question_match = re.search(
        r"USER QUESTION / OBJECTIVE:\n(.*?)\n\nALLOWED CITATION", prompt, re.DOTALL
    )
    question_text = question_match.group(1).strip() if question_match else ""
    question = next(
        (
            item
            for item in cast(Sequence[Mapping[str, object]], benchmark["questions"])
            if item["question"] == question_text
        ),
        None,
    )
    if question is None:
        payload = {
            "schema_version": SCHEMA_VERSION,
            "paragraphs": [],
            "claims": [],
            "insufficient_evidence": ["No deterministic answer is defined."],
            "contradictions": [],
        }
        return json.dumps(payload)
    _, by_filename = _source_maps(benchmark)
    citations: dict[str, str] = {}
    for citation_id, filename in re.findall(
        r"^\[(CIT-\d+)\] Source: (.+?) \(", prompt, re.MULTILINE
    ):
        source_id = by_filename.get(filename)
        if source_id:
            citations[source_id] = citation_id
    expected = [str(item) for item in cast(Sequence[object], question["expected_sources"])]
    if not expected:
        payload = {
            "schema_version": SCHEMA_VERSION,
            "paragraphs": [],
            "claims": [],
            "insufficient_evidence": ["The available evidence does not establish this."],
            "contradictions": [],
        }
        return json.dumps(payload)
    if any(source_id not in citations for source_id in expected):
        payload = {
            "schema_version": SCHEMA_VERSION,
            "paragraphs": [],
            "claims": [],
            "insufficient_evidence": ["Required evidence was absent from the retrieval package."],
            "contradictions": [],
        }
        return json.dumps(payload)
    facts = cast(Sequence[Mapping[str, object]], question["required_facts"])
    text = (
        "; ".join(
            " ".join(str(term) for term in cast(Sequence[object], fact["terms"]))
            for fact in facts
        )
        or "The requested information is supported."
    )
    payload = {
        "schema_version": SCHEMA_VERSION,
        "paragraphs": [
            {
                "paragraph_id": "PAR-1",
                "text": text,
                "citation_ids": [citations[source_id] for source_id in expected],
            }
        ],
        "claims": [],
        "insufficient_evidence": [],
        "contradictions": (
            ["The retrieved sources contain materially conflicting statements."]
            if question["expected_contradiction"]
            else []
        ),
    }
    return json.dumps(payload)


class FakeModelServer:
    """Deterministic loopback OpenAI-compatible server for public CI only."""

    def __init__(self, benchmark: Mapping[str, object], mode: str = "valid") -> None:
        self.benchmark = benchmark
        self.mode = mode
        self.server: ThreadingHTTPServer | None = None
        self.thread: threading.Thread | None = None

    def __enter__(self) -> "FakeModelServer":
        benchmark = self.benchmark
        mode = self.mode

        class Handler(BaseHTTPRequestHandler):
            def do_POST(self) -> None:  # noqa: N802
                if self.path != "/v1/chat/completions":
                    self.send_error(404)
                    return
                length = int(self.headers.get("Content-Length", "0"))
                body = json.loads(self.rfile.read(length).decode("utf-8"))
                prompt = str(body["messages"][0]["content"])
                if mode == "error":
                    self.send_error(503, "synthetic failure")
                    return
                content = "not-json" if mode == "invalid" else _fake_response(benchmark, prompt)
                encoded = json.dumps({"choices": [{"message": {"content": content}}]}).encode(
                    "utf-8"
                )
                self.send_response(200)
                self.send_header("Content-Type", "application/json")
                self.send_header("Content-Length", str(len(encoded)))
                self.end_headers()
                self.wfile.write(encoded)

            def log_message(self, format: str, *args: object) -> None:
                del format, args

        self.server = ThreadingHTTPServer(("127.0.0.1", 0), Handler)
        self.thread = threading.Thread(target=self.server.serve_forever, daemon=True)
        self.thread.start()
        return self

    def __exit__(self, *args: object) -> None:
        del args
        if self.server:
            self.server.shutdown()
            self.server.server_close()
        if self.thread:
            self.thread.join(timeout=5)

    @property
    def endpoint(self) -> str:
        if not self.server:
            raise RuntimeError("fake model server is not running")
        return f"http://127.0.0.1:{self.server.server_port}"
