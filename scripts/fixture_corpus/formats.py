"""Deterministic builders for the binary fixture formats."""

from __future__ import annotations

import math
import re
import struct
from io import BytesIO
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

from docx import Document
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from reportlab.pdfgen.canvas import Canvas

from fixture_corpus.specs import FIXED_DATETIME, FIXED_TIME


def _normalize_zip(raw: bytes) -> bytes:
    source = BytesIO(raw)
    target = BytesIO()
    with ZipFile(source) as input_zip, ZipFile(target, "w") as output_zip:
        for name in sorted(input_zip.namelist()):
            info = ZipInfo(name, FIXED_TIME)
            info.compress_type = ZIP_DEFLATED
            info.create_system = 3
            info.external_attr = 0o100644 << 16
            content = input_zip.read(name)
            if name == "docProps/core.xml":
                content = re.sub(
                    rb"<dcterms:(created|modified)([^>]*)>[^<]+</dcterms:\1>",
                    lambda match: (
                        b"<dcterms:"
                        + match.group(1)
                        + match.group(2)
                        + b">2026-01-01T00:00:00Z</dcterms:"
                        + match.group(1)
                        + b">"
                    ),
                    content,
                )
            output_zip.writestr(info, content)
    return target.getvalue()


def build_docx() -> bytes:
    document = Document()
    document.core_properties.title = "Archiv DOCX Fixture"
    document.core_properties.author = "Archiv fixture generator"
    document.core_properties.created = FIXED_DATETIME
    document.core_properties.modified = FIXED_DATETIME
    document.add_paragraph("Archiv DOCX Fixture")
    document.add_paragraph("ARCHIV-DOCX-MARKER-2026")
    document.add_paragraph("Section: Findings; paragraph location: 2")
    raw = BytesIO()
    document.save(raw)
    return _normalize_zip(raw.getvalue())


def build_xlsx() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Evidence"
    sheet["A1"] = "Archiv XLSX Fixture"
    sheet["B2"] = "ARCHIV-XLSX-MARKER-2026"
    sheet["C3"] = 42.25
    workbook.properties.title = "Archiv XLSX Fixture"
    workbook.properties.creator = "Archiv fixture generator"
    workbook.properties.created = FIXED_DATETIME
    workbook.properties.modified = FIXED_DATETIME
    raw = BytesIO()
    workbook.save(raw)
    return _normalize_zip(raw.getvalue())


def build_pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(457200, 457200, 8229600, 914400)
    title.text_frame.text = "Archiv PPTX Fixture"
    marker = slide.shapes.add_textbox(457200, 1828800, 8229600, 914400)
    marker.text_frame.text = "ARCHIV-PPTX-MARKER-2026"
    presentation.core_properties.title = "Archiv PPTX Fixture"
    presentation.core_properties.author = "Archiv fixture generator"
    presentation.core_properties.created = FIXED_DATETIME
    presentation.core_properties.modified = FIXED_DATETIME
    raw = BytesIO()
    presentation.save(raw)
    return _normalize_zip(raw.getvalue())


def build_pdf() -> bytes:
    raw = BytesIO()
    canvas = Canvas(
        raw,
        pagesize=(612, 792),
        invariant=1,
        pageCompression=0,
        pdfVersion=(1, 4),
    )
    canvas.setTitle("Archiv PDF Fixture")
    canvas.setAuthor("Archiv fixture generator")
    canvas.setFont("Helvetica", 18)
    canvas.drawString(72, 720, "Archiv PDF Fixture")
    canvas.drawString(72, 692, "ARCHIV-PDF-MARKER-2026")
    canvas.drawString(72, 664, "Page location: 1")
    canvas.showPage()
    canvas.save()
    return raw.getvalue()


def build_png() -> bytes:
    image = Image.new("L", (640, 80), 255)
    draw = ImageDraw.Draw(image)
    draw.text(
        (20, 20),
        "ARCHIV-IMAGE-MARKER-2026",
        fill=0,
        font=ImageFont.load_default(),
    )
    raw = BytesIO()
    image.save(raw, format="PNG", optimize=False, compress_level=9)
    return raw.getvalue()


def _chunk(kind: bytes, payload: bytes) -> bytes:
    return kind + struct.pack("<I", len(payload)) + payload + (b"\0" if len(payload) % 2 else b"")


def build_wav() -> bytes:
    sample_rate = 8000
    tones = (
        (440.0, 0.25),
        (0.0, 0.05),
        (554.37, 0.25),
        (0.0, 0.05),
        (659.25, 0.25),
    )
    samples: list[int] = []
    for frequency, duration in tones:
        for index in range(round(sample_rate * duration)):
            value = (
                0
                if frequency == 0
                else int(9000 * math.sin(2 * math.pi * frequency * index / sample_rate))
            )
            samples.append(value)
    pcm = struct.pack("<" + "h" * len(samples), *samples)
    fmt = struct.pack("<HHIIHH", 1, 1, sample_rate, sample_rate * 2, 2, 16)
    info = b"INFO" + _chunk(b"ICMT", b"ARCHIV-AUDIO-MARKER-2026\0")
    chunks = _chunk(b"fmt ", fmt) + _chunk(b"data", pcm) + _chunk(b"LIST", info)
    return b"RIFF" + struct.pack("<I", len(chunks) + 4) + b"WAVE" + chunks
